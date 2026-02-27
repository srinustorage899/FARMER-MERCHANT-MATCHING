"""
Generate the AgriMatch Project Review PPT.
Matches the exact structure of the sample PPT (13 slides).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ───────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)   # 16:9
SLIDE_HEIGHT = Inches(7.5)

# Colour palette
BG_DARK     = RGBColor(0x1B, 0x1B, 0x2F)   # dark navy
BG_CARD     = RGBColor(0x24, 0x24, 0x3E)   # card bg
ACCENT      = RGBColor(0x00, 0xC9, 0xA7)   # teal/green accent
ACCENT2     = RGBColor(0x4F, 0x9D, 0xF7)   # blue accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MUTED       = RGBColor(0x99, 0x99, 0xAA)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
RED         = RGBColor(0xFF, 0x45, 0x45)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
YELLOW      = RGBColor(0xFF, 0xD6, 0x00)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helpers ─────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=16, bold=False, color=WHITE,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri",
             bullet=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def add_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=13):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=11, color=MUTED,
                alignment=PP_ALIGN.RIGHT)


def add_footer_bar(slide, num, total=13):
    add_rect(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.3), ACCENT)
    add_slide_number(slide, num, total)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Accent strip at top
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

# Student info - left
add_textbox(slide, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4),
            "160122733305 – POTHARAJU SRINIVAS", font_size=14, color=LIGHT_GRAY)
add_textbox(slide, Inches(0.6), Inches(0.75), Inches(6), Inches(0.4),
            "160122733302 – ENDRAVATH KRISHNA", font_size=14, color=LIGHT_GRAY)

# Supervisor - right
add_textbox(slide, Inches(8.5), Inches(0.4), Inches(4.5), Inches(0.4),
            "Supervisor", font_size=12, color=MUTED, alignment=PP_ALIGN.RIGHT)
add_textbox(slide, Inches(8.5), Inches(0.7), Inches(4.5), Inches(0.4),
            "DR. M SWAMY DAS", font_size=14, bold=True, color=ACCENT,
            alignment=PP_ALIGN.RIGHT)

# PID
add_textbox(slide, Inches(0.6), Inches(1.2), Inches(3), Inches(0.4),
            "PID - 36", font_size=14, bold=True, color=ORANGE)

# Main title
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.2),
            "AgriMatch: HNSW-Powered Farmer-Merchant\nMatching Platform",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
            font_name="Calibri Light")

# Subtitle
add_textbox(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.8),
            "A high-performance agricultural marketplace using Hierarchical Navigable Small World "
            "graphs for real-time vector similarity matching between farmers and merchants",
            font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Accent line under title
add_accent_line(slide, Inches(5.5), Inches(4.5), Inches(2.3))

# Tech badges at bottom
badges = ["FastAPI", "React.js", "MongoDB", "Redis", "HNSW (hnswlib)", "Docker"]
badge_x_start = Inches(2.2)
for i, badge in enumerate(badges):
    x = badge_x_start + Inches(i * 1.6)
    r = add_rect(slide, x, Inches(5.8), Inches(1.4), Inches(0.38), BG_CARD, ACCENT)
    add_textbox(slide, x, Inches(5.8), Inches(1.4), Inches(0.38),
                badge, font_size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide, 1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
            "Problem Statement", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

# Problem text in a card
card = add_rect(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(5.4), BG_CARD, RGBColor(0x33,0x33,0x55))
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = ("Indian agriculture remains heavily dependent on intermediary-driven supply "
          "chains, where farmers receive only a fraction of the final selling price. "
          "The absence of a direct digital platform connecting farmers with merchants "
          "leads to significant price inefficiencies — farmers cannot access fair market "
          "rates, while merchants face inflated procurement costs.")
p.font.size = Pt(17)
p.font.color.rgb = LIGHT_GRAY
p.font.name = "Calibri"
p.line_spacing = Pt(28)

add_para(tf, "", font_size=8)

add_para(tf, ("Furthermore, existing agricultural marketplaces lack intelligent matching "
              "capabilities that consider multiple factors such as crop type, quantity, "
              "pricing, quality, and geographical proximity simultaneously. This gap "
              "results in suboptimal trade pairings, increased transportation costs, "
              "and missed opportunities for both parties."),
         font_size=17, color=LIGHT_GRAY)

add_para(tf, "", font_size=8)

add_para(tf, ("Moreover, the lack of integration with government-mandated Minimum Support "
              "Prices (MSP) means farmers often sell below fair rates without any system-level "
              "safeguard. The current ecosystem offers no mechanism for real-time, "
              "proximity-aware matching that could optimize logistics and minimize wastage — "
              "contributing to large-scale economic loss and rural livelihood challenges."),
         font_size=17, color=LIGHT_GRAY)

add_footer_bar(slide, 2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Design (System Flow)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
            "Design — System Flow", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

# Farmer Flow (Left side)
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
            "🌾  Farmer Upload Flow", font_size=20, bold=True, color=ACCENT)

farmer_steps = [
    ("1", "Farmer opens Upload Crop page"),
    ("2", "Enters crop, quantity, price + GPS auto-detected"),
    ("3", "POST /api/farmer/upload → FastAPI"),
    ("4", "Validated via Pydantic schemas"),
    ("5", "Stored in MongoDB (listings collection)"),
    ("6", "Listing ID pushed to Redis queue"),
    ("7", "Background worker encodes 54-D vector"),
    ("8", "Inserted into per-crop HNSW index"),
]

for i, (num, text) in enumerate(farmer_steps):
    y = Inches(2.0) + Inches(i * 0.55)
    circle = add_rect(slide, Inches(1.0), y, Inches(0.35), Inches(0.35), ACCENT)
    add_textbox(slide, Inches(1.0), y, Inches(0.35), Inches(0.35),
                num, font_size=12, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), y, Inches(5), Inches(0.35),
                text, font_size=13, color=LIGHT_GRAY)

# Arrow in middle
add_textbox(slide, Inches(6.2), Inches(3.5), Inches(0.8), Inches(1),
            "→", font_size=48, color=ACCENT, alignment=PP_ALIGN.CENTER)

# Merchant Flow (Right side)
add_textbox(slide, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
            "🔍  Merchant Search Flow", font_size=20, bold=True, color=ACCENT2)

merchant_steps = [
    ("1", "Merchant opens Find Farmers page"),
    ("2", "Enters crop, qty, max price, radius + GPS"),
    ("3", "POST /api/merchant/search → FastAPI"),
    ("4", "Query encoded into same 54-D vector space"),
    ("5", "HNSW ANN search → top-50 candidates"),
    ("6", "Filter: active + MSP + budget constraints"),
    ("7", "Haversine reranking by real distance"),
    ("8", "Return top-10 matches sorted by proximity"),
]

for i, (num, text) in enumerate(merchant_steps):
    y = Inches(2.0) + Inches(i * 0.55)
    circle = add_rect(slide, Inches(7.4), y, Inches(0.35), Inches(0.35), ACCENT2)
    add_textbox(slide, Inches(7.4), y, Inches(0.35), Inches(0.35),
                num, font_size=12, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.9), y, Inches(5), Inches(0.35),
                text, font_size=13, color=LIGHT_GRAY)

add_footer_bar(slide, 3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Design (Data Flow / Vector Pipeline)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
            "Design — Vector Encoding Pipeline", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

# Vector layout diagram
card1 = add_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5), BG_CARD, RGBColor(0x33,0x33,0x55))

add_textbox(slide, Inches(0.9), Inches(1.5), Inches(5.5), Inches(0.5),
            "54-Dimensional Feature Vector", font_size=20, bold=True, color=ACCENT)

# Vector components
components = [
    ("Dim 0", "Latitude (normalised 0–1)", "GPS latitude mapped to [0,1]"),
    ("Dim 1", "Longitude (normalised 0–1)", "GPS longitude mapped to [0,1]"),
    ("Dim 2", "Price (normalised 0–1)", "₹/kg mapped via [0, 500] range"),
    ("Dim 3", "Quantity (normalised 0–1)", "kg mapped via [0, 50000] range"),
    ("Dim 4–53", "Crop One-Hot (50 dims)", "50 supported crops encoded"),
]

for i, (dim, name, desc) in enumerate(components):
    y = Inches(2.2) + Inches(i * 0.8)
    add_rect(slide, Inches(1.0), y, Inches(1.2), Inches(0.55), ACCENT if i < 4 else ACCENT2)
    add_textbox(slide, Inches(1.0), y, Inches(1.2), Inches(0.55),
                dim, font_size=11, bold=True, color=BG_DARK, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.4), y, Inches(3.5), Inches(0.3),
                name, font_size=14, bold=True, color=WHITE)
    add_textbox(slide, Inches(2.4), Inches(0.28) + y, Inches(3.5), Inches(0.3),
                desc, font_size=11, color=MUTED)

# Normalisation formula
add_textbox(slide, Inches(0.9), Inches(6.2), Inches(5.5), Inches(0.4),
            "Normalisation: x' = (x − min) / (max − min),  clamped to [0, 1]",
            font_size=12, color=MUTED)

# Right side — HNSW Parameters
card2 = add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5), BG_CARD, RGBColor(0x33,0x33,0x55))

add_textbox(slide, Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.5),
            "HNSW Index Configuration", font_size=20, bold=True, color=ACCENT)

params = [
    ("Space", "L2 (Euclidean)", "Euclidean distance between normalised vectors"),
    ("M", "16", "Max bi-directional links per node per layer"),
    ("ef_construction", "200", "Candidate list size during index build"),
    ("ef_search", "100", "Candidate list size during query"),
    ("Max Elements", "10,000", "Initial capacity (auto-resized)"),
    ("Search K", "50", "Top-K candidates before reranking"),
    ("Indexing", "Per-crop", "Separate HNSW index for each crop type"),
    ("Persistence", ".hnsw + .meta.json", "Saved to disk on shutdown, loaded on startup"),
]

for i, (key, val, desc) in enumerate(params):
    y = Inches(2.2) + Inches(i * 0.58)
    add_textbox(slide, Inches(7.2), y, Inches(2.5), Inches(0.3),
                key, font_size=13, bold=True, color=ACCENT2)
    add_textbox(slide, Inches(9.5), y, Inches(1.2), Inches(0.3),
                val, font_size=13, bold=True, color=YELLOW)
    add_textbox(slide, Inches(7.2), Inches(0.25) + y, Inches(5), Inches(0.3),
                desc, font_size=10, color=MUTED)

add_footer_bar(slide, 4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
            "Architecture", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

# Architecture layers as horizontal bands

layers = [
    ("Frontend Layer", "React.js + Vite  |  Responsive SPA  |  Geolocation API  |  Multi-language (EN/HI/TE)",
     ACCENT, Inches(1.4)),
    ("API Layer", "FastAPI (Python)  |  Pydantic Validation  |  RESTful Endpoints  |  CORS Middleware",
     ACCENT2, Inches(2.5)),
    ("Service Layer", "Vector Service  |  Feature Engineering  |  MSP Validation  |  Haversine Reranking",
     ORANGE, Inches(3.6)),
    ("Data Layer", "MongoDB (Documents)  |  Redis (Job Queue)  |  HNSW Indices (Per-Crop .hnsw files)",
     RGBColor(0xAB,0x47,0xBC), Inches(4.7)),
    ("Infrastructure", "Docker Compose  |  Background Worker  |  Auto-reload Dev Server  |  Index Persistence",
     RGBColor(0xEF,0x53,0x50), Inches(5.8)),
]

for title, desc, color, y in layers:
    bar = add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), BG_CARD, color)
    add_textbox(slide, Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.4),
                title, font_size=16, bold=True, color=color)
    add_textbox(slide, Inches(1.0), y + Inches(0.42), Inches(11), Inches(0.4),
                desc, font_size=13, color=LIGHT_GRAY)

# Arrows between layers
for i in range(4):
    y = Inches(2.25) + Inches(i * 1.1)
    add_textbox(slide, Inches(6.2), y, Inches(1), Inches(0.35),
                "▼", font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide, 5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Implementation
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
            "Implementation", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

# Implementation cards — 2 columns × 3 rows
impl_items = [
    ("Vector Search Engine", ACCENT,
     "• 54-D feature vectors (4 numeric + 50 crop one-hot)\n"
     "• Per-crop HNSW indices via hnswlib\n"
     "• Two-phase: ANN search → Haversine rerank\n"
     "• Sub-linear O(log n) search complexity"),
    ("MSP Price Validation", GREEN,
     "• Government MSP rates for all 50 crops\n"
     "• Server-side enforcement — no below-MSP listings\n"
     "• Protects farmers from unfair pricing\n"
     "• Auto-validated on every search result"),
    ("Async Queue Pipeline", ACCENT2,
     "• Redis-backed BRPOP insertion queue\n"
     "• Background worker for HNSW indexing\n"
     "• Non-blocking upload → instant farmer response\n"
     "• Decouples API from heavy index operations"),
    ("Geolocation & Distance", ORANGE,
     "• Browser Geolocation API for auto GPS\n"
     "• Haversine great-circle distance (km)\n"
     "• Radius-based filtering post-search\n"
     "• Results sorted by real-world proximity"),
    ("Data Persistence", RGBColor(0xAB,0x47,0xBC),
     "• HNSW indices saved as .hnsw binary files\n"
     "• Metadata in .meta.json (ID mapping)\n"
     "• Auto-load on startup, save on shutdown\n"
     "• MongoDB for all listing documents"),
    ("Frontend SPA", RGBColor(0xEF, 0x53, 0x50),
     "• React.js with Vite build system\n"
     "• Protected routes + auth context\n"
     "• Multi-language support (EN/HI/TE)\n"
     "• Responsive dashboard for both roles"),
]

for i, (title, color, desc) in enumerate(impl_items):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(1.3) + Inches(row * 2.0)

    card = add_rect(slide, x, y, Inches(5.9), Inches(1.8), BG_CARD, color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.4),
                title, font_size=16, bold=True, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(1.3),
                desc, font_size=12, color=LIGHT_GRAY)

add_footer_bar(slide, 6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDES 7–11 — Application Screenshots (placeholder structure)
# ═══════════════════════════════════════════════════════════════════════════
screenshot_slides = [
    ("Homepage & Authentication", "Landing page with role-based login (Farmer / Merchant)"),
    ("Farmer Dashboard & Upload", "Crop upload form with GPS auto-detection and MSP display"),
    ("Merchant Search & Results", "Find Farmers page with HNSW-powered matching results"),
    ("Farmers Database (Live)", "Real-time MongoDB viewer with 100+ listings and location names"),
    ("HNSW Visualization (3D)", "Interactive Three.js visualization of HNSW graph traversal"),
]

for i, (title, desc) in enumerate(screenshot_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                title, font_size=28, bold=True, color=WHITE)
    add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(2.5))

    # Placeholder box for screenshot
    placeholder = add_rect(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.6),
                           RGBColor(0x1A, 0x1A, 0x2A), RGBColor(0x44, 0x44, 0x66))
    add_textbox(slide, Inches(3.5), Inches(3.2), Inches(6), Inches(1.5),
                f"[ Insert Screenshot Here ]\n{desc}",
                font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

    add_footer_bar(slide, 7 + i)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Links / References
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6),
            "References & Links", font_size=32, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2.5))

# Link cards
links = [
    ("Technologies Used", ACCENT, [
        ("Backend", "FastAPI (Python 3.11) — High-performance async web framework"),
        ("Frontend", "React.js 18 + Vite — Modern SPA with hot module reload"),
        ("Database", "MongoDB 7 — Document store for farmer listings"),
        ("Queue", "Redis 7 — In-memory job queue for async HNSW insertion"),
        ("Vector Search", "hnswlib — C++ HNSW with Python bindings"),
        ("Visualization", "Three.js — WebGL 3D HNSW graph visualization"),
        ("Containerization", "Docker Compose — Multi-service orchestration"),
    ]),
    ("Key References", ACCENT2, [
        ("HNSW Paper", "Malkov & Yashunin, 'Efficient and robust approximate nearest neighbor using HNSW graphs', IEEE TPAMI 2018"),
        ("MSP Data", "Government of India — Ministry of Agriculture & Farmers' Welfare, MSP rates 2024-25"),
        ("Haversine", "Great-circle distance formula for geospatial proximity calculation"),
        ("Nominatim", "OpenStreetMap reverse geocoding API for human-readable location names"),
    ]),
]

y_offset = Inches(1.4)
for section_title, color, items in links:
    add_textbox(slide, Inches(0.8), y_offset, Inches(10), Inches(0.5),
                section_title, font_size=20, bold=True, color=color)
    y_offset += Inches(0.5)

    for key, val in items:
        add_textbox(slide, Inches(1.2), y_offset, Inches(2), Inches(0.35),
                    key, font_size=13, bold=True, color=WHITE)
        add_textbox(slide, Inches(3.5), y_offset, Inches(9), Inches(0.35),
                    val, font_size=12, color=LIGHT_GRAY)
        y_offset += Inches(0.35)

    y_offset += Inches(0.3)

add_footer_bar(slide, 12)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5),
            "THANK YOU!", font_size=56, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

add_accent_line(slide, Inches(5.5), Inches(3.7), Inches(2.3))

add_textbox(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.6),
            "AgriMatch — Connecting Farmers & Merchants with AI-Powered Matching",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(5.2), Inches(9), Inches(0.4),
            "POTHARAJU SRINIVAS  •  ENDRAVATH KRISHNA",
            font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
            "Supervisor: DR. M SWAMY DAS  |  PID-36",
            font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide, 13)


# ─── Save ───────────────────────────────────────────────────────────────────
output_path = r"d:\HOPE\AgriMatch_Project_Review_PID36.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
